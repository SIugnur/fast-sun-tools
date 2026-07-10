import os
from PyQt5.QtWidgets import (QWidget, QVBoxLayout, QLabel, QScrollArea, 
                             QTextEdit, QFrame, QSizePolicy, QTableWidget,
                             QTableWidgetItem, QHeaderView)
from PyQt5.QtCore import Qt, QSize
from PyQt5.QtGui import QPixmap, QImage, QFont, QMovie
from PIL import Image


class FilePreview(QWidget):
    def __init__(self):
        super().__init__()
        self.init_ui()
        
    def init_ui(self):
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        
        self.title_label = QLabel("文件预览")
        self.title_label.setStyleSheet("font-weight: bold; padding: 8px; background: #f0f0f0;")
        layout.addWidget(self.title_label)
        
        self.scroll_area = QScrollArea()
        self.scroll_area.setWidgetResizable(True)
        self.scroll_area.setHorizontalScrollBarPolicy(Qt.ScrollBarAsNeeded)
        self.scroll_area.setVerticalScrollBarPolicy(Qt.ScrollBarAsNeeded)
        
        self.preview_content = QLabel("选择一个文件进行预览")
        self.preview_content.setAlignment(Qt.AlignCenter)
        self.preview_content.setStyleSheet("color: #888; font-size: 14px;")
        self.scroll_area.setWidget(self.preview_content)
        
        layout.addWidget(self.scroll_area)
        
    def preview_file(self, file_path):
        if not os.path.isfile(file_path):
            self.show_message("这是一个文件夹")
            return
            
        ext = os.path.splitext(file_path)[1].lower()
        self.title_label.setText(f"预览: {os.path.basename(file_path)}")
        
        if ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff']:
            self.preview_image(file_path)
        elif ext in ['.txt', '.py', '.js', '.html', '.css', '.json', '.xml', '.md', '.csv', '.log', '.ini', '.cfg']:
            self.preview_text(file_path)
        elif ext in ['.pdf']:
            self.preview_pdf(file_path)
        elif ext in ['.docx', '.doc']:
            self.preview_word(file_path)
        elif ext in ['.xlsx', '.xls']:
            self.preview_excel(file_path)
        elif ext in ['.pptx', '.ppt']:
            self.preview_powerpoint(file_path)
        elif ext in ['.mp4', '.avi', '.mov', '.mkv', '.wmv']:
            self.show_message("视频文件: 请使用视频播放器打开")
        elif ext in ['.mp3', '.wav', '.flac', '.aac']:
            self.show_message("音频文件: 请使用音频播放器打开")
        else:
            self.show_message(f"不支持的文件类型: {ext}")
            
    def preview_image(self, file_path):
        try:
            pixmap = QPixmap(file_path)
            if not pixmap.isNull():
                label = QLabel()
                scaled = pixmap.scaled(
                    self.scroll_area.size() - QSize(20, 20),
                    Qt.KeepAspectRatio,
                    Qt.SmoothTransformation
                )
                label.setPixmap(scaled)
                label.setAlignment(Qt.AlignCenter)
                self.scroll_area.setWidget(label)
            else:
                self.show_message("无法加载图片")
        except Exception as e:
            self.show_message(f"预览图片失败: {str(e)}")
            
    def preview_text(self, file_path):
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            text_edit = QTextEdit()
            text_edit.setReadOnly(True)
            text_edit.setFont(QFont("Consolas", 10))
            text_edit.setPlainText(content)
            self.scroll_area.setWidget(text_edit)
        except Exception as e:
            self.show_message(f"预览文本失败: {str(e)}")
            
    def preview_pdf(self, file_path):
        try:
            import fitz  # PyMuPDF
            
            # 使用 PyMuPDF 打开 PDF
            doc = fitz.open(file_path)
            num_pages = len(doc)
            
            if num_pages == 0:
                self.show_message("PDF 文档为空")
                return
            
            # 创建容器widget
            container = QWidget()
            container_layout = QVBoxLayout(container)
            container_layout.setContentsMargins(10, 10, 10, 10)
            
            # 信息标签
            info_label = QLabel(f"PDF 文档预览 (共 {num_pages} 页)")
            info_label.setStyleSheet("font-weight: bold; font-size: 14px; padding: 5px;")
            container_layout.addWidget(info_label)
            
            # 显示前 5 页的缩略图
            max_preview_pages = min(5, num_pages)
            
            for page_num in range(max_preview_pages):
                page = doc[page_num]
                
                # 渲染页面为图片（高分辨率）
                zoom = 2.0  # 缩放因子，提高清晰度
                mat = fitz.Matrix(zoom, zoom)
                pix = page.get_pixmap(matrix=mat)
                
                # 转换为 Qt 图片
                img_data = pix.tobytes("png")
                qimage = QImage.fromData(img_data)
                pixmap = QPixmap.fromImage(qimage)
                
                # 页面标签
                page_label = QLabel(f"--- 第 {page_num + 1} 页 ---")
                page_label.setStyleSheet("font-weight: bold; color: #2196F3; padding: 5px 0;")
                container_layout.addWidget(page_label)
                
                # 图片标签
                img_label = QLabel()
                img_label.setPixmap(pixmap)
                img_label.setAlignment(Qt.AlignCenter)
                container_layout.addWidget(img_label)
                
                # 尝试提取文本（如果存在）
                text = page.get_text()
                if text.strip():
                    text_label = QLabel("📝 文本内容:")
                    text_label.setStyleSheet("font-weight: bold; color: #666;")
                    container_layout.addWidget(text_label)
                    
                    text_edit = QTextEdit()
                    text_edit.setReadOnly(True)
                    text_edit.setMaximumHeight(150)
                    text_edit.setFont(QFont("Microsoft YaHei", 9))
                    text_edit.setPlainText(text[:500] + "..." if len(text) > 500 else text)
                    container_layout.addWidget(text_edit)
                
                container_layout.addWidget(QLabel())  # 间距
            
            doc.close()
            
            if max_preview_pages < num_pages:
                more_label = QLabel(f"⚠️ 仅显示前 {max_preview_pages} 页，共 {num_pages} 页")
                more_label.setStyleSheet("color: #ff9800; font-style: italic; padding: 10px;")
                container_layout.addWidget(more_label)
            
            self.scroll_area.setWidget(container)
            
        except ImportError:
            self.show_message("PDF 预览: 需要安装 PyMuPDF 库\n请运行: pip install PyMuPDF")
        except Exception as e:
            self.show_message(f"PDF 预览失败: {str(e)}\n\n提示：如果是扫描版PDF，可能无法预览文本内容")
            
    def preview_word(self, file_path):
        try:
            from docx import Document
            
            doc = Document(file_path)
            paragraphs = []
            
            # 提取段落文本（前 100 个段落）
            for i, para in enumerate(doc.paragraphs[:100]):
                if para.text.strip():
                    paragraphs.append(para.text)
            
            # 提取表格内容
            tables_content = []
            for table in doc.tables[:5]:  # 最多 5 个表格
                table_text = []
                for row in table.rows[:10]:  # 每个表格最多 10 行
                    cells = [cell.text.strip() for cell in row.cells]
                    table_text.append(" | ".join(cells))
                if table_text:
                    tables_content.append("\n".join(table_text))
            
            content = f"Word 文档预览\n{'='*50}\n\n"
            content += "【段落内容】\n"
            content += "\n".join(paragraphs) if paragraphs else "无段落内容"
            
            if tables_content:
                content += f"\n\n【表格内容】\n"
                content += "\n\n".join(tables_content)
            
            text_edit = QTextEdit()
            text_edit.setReadOnly(True)
            text_edit.setFont(QFont("Microsoft YaHei", 10))
            text_edit.setPlainText(content)
            self.scroll_area.setWidget(text_edit)
            
        except ImportError:
            self.show_message("Word 预览: 需要安装 python-docx 库\n请运行: pip install python-docx")
        except Exception as e:
            self.show_message(f"Word 预览失败: {str(e)}")
            
    def preview_excel(self, file_path):
        try:
            import openpyxl
            from openpyxl.utils.exceptions import InvalidFileException
            
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            
            # 预览第一个工作表
            sheet_names = wb.sheetnames
            first_sheet_name = sheet_names[0] if sheet_names else "Sheet1"
            ws = wb[first_sheet_name]
            
            # 创建表格widget
            table_widget = QTableWidget()
            
            # 设置行数和列数（最多显示 50 行 20 列）
            max_rows = min(50, ws.max_row or 1)
            max_cols = min(20, ws.max_column or 1)
            
            table_widget.setRowCount(max_rows)
            table_widget.setColumnCount(max_cols)
            
            # 设置表头
            headers = []
            for col in range(max_cols):
                cell_value = ws.cell(row=1, column=col+1).value
                headers.append(str(cell_value) if cell_value else f"列{col+1}")
            table_widget.setHorizontalHeaderLabels(headers)
            
            # 填充数据
            for row in range(max_rows):
                for col in range(max_cols):
                    cell_value = ws.cell(row=row+1, column=col+1).value
                    item = QTableWidgetItem(str(cell_value) if cell_value else "")
                    table_widget.setItem(row, col, item)
            
            # 调整列宽
            table_widget.horizontalHeader().setSectionResizeMode(QHeaderView.Stretch)
            
            # 添加说明标签
            info_label = QLabel(f"Excel 工作表预览: {first_sheet_name} (共 {ws.max_row or 0} 行 x {ws.max_column or 0} 列)")
            info_label.setStyleSheet("font-weight: bold; padding: 8px; background: #e3f2fd;")
            
            container = QWidget()
            layout = QVBoxLayout(container)
            layout.setContentsMargins(0, 0, 0, 0)
            layout.addWidget(info_label)
            layout.addWidget(table_widget)
            
            self.scroll_area.setWidget(container)
            wb.close()
            
        except ImportError:
            self.show_message("Excel 预览: 需要安装 openpyxl 库\n请运行: pip install openpyxl")
        except Exception as e:
            self.show_message(f"Excel 预览失败: {str(e)}")
            
    def preview_powerpoint(self, file_path):
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            
            slides_content = []
            
            # 提取前 10 张幻灯片的文本
            for i, slide in enumerate(prs.slides[:10]):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            slide_text.append(shape.text.strip())
                
                if slide_text:
                    slides_content.append(f"【幻灯片 {i+1}】\n" + "\n".join(slide_text[:10]))
            
            content = f"PowerPoint 预览 (共 {len(prs.slides)} 张幻灯片)\n"
            content += f"显示前 {len(slides_content)} 张内容：\n"
            content += f"\n{'='*50}\n\n"
            content += "\n\n".join(slides_content) if slides_content else "无法提取文本内容"
            
            text_edit = QTextEdit()
            text_edit.setReadOnly(True)
            text_edit.setFont(QFont("Microsoft YaHei", 10))
            text_edit.setPlainText(content)
            self.scroll_area.setWidget(text_edit)
            
        except ImportError:
            self.show_message("PowerPoint 预览: 需要安装 python-pptx 库\n请运行: pip install python-pptx")
        except Exception as e:
            self.show_message(f"PowerPoint 预览失败: {str(e)}")
            
    def show_message(self, message):
        label = QLabel(message)
        label.setAlignment(Qt.AlignCenter)
        label.setStyleSheet("color: #666; font-size: 13px; padding: 20px;")
        label.setWordWrap(True)
        self.scroll_area.setWidget(label)
